import os
import sys
import subprocess

# Auto-install python-pptx if it is not present
try:
    import pptx
except ImportError:
    print("[AURA] python-pptx not found. Installing now...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_pitch_deck():
    prs = Presentation()
    
    # Set slide dimensions to 16:9 Widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Color palette definitions
    NAVY = RGBColor(14, 30, 53)        # #0E1E35 (Primary background)
    DARK_NAVY = RGBColor(7, 15, 30)    # #070F1E (Lighter/Darker highlights)
    ORANGE = RGBColor(200, 82, 26)     # #C8521A (Accent color)
    LIGHT_ORANGE = RGBColor(255, 138, 80) # #FF8A50 (Muted highlight)
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(148, 163, 184)     # #94A3B8 (Body text)
    RED = RGBColor(239, 68, 68)        # #EF4444
    GREEN = RGBColor(16, 185, 129)     # #10B981
    
    # Helper to set solid background color
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper to add a header zone
    def add_header(slide, title_text, subtitle_text=""):
        # Header shape
        txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(1.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = 'Arial'
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.size = Pt(14)
            p2.font.color.rgb = LIGHT_ORANGE
            p2.font.name = 'Arial'
            p2.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 1: Title / Intro Slide
    # -------------------------------------------------------------
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide1 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide1, NAVY)
    
    # Design accents (orange stripe on the left)
    accent_bar = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ORANGE
    accent_bar.line.fill.background() # No line
    
    # Logo text
    logo_box = slide1.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(4), Inches(0.5))
    ltf = logo_box.text_frame
    lp = ltf.paragraphs[0]
    lp.text = "AURA  ▪"
    lp.font.size = Pt(14)
    lp.font.bold = True
    lp.font.color.rgb = LIGHT_ORANGE
    
    # Main Title Box
    title_box = slide1.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(11.5), Inches(3.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p_title = tf.paragraphs[0]
    p_title.text = "AURA"
    p_title.font.size = Pt(64)
    p_title.font.bold = True
    p_title.font.color.rgb = WHITE
    p_title.font.name = 'Arial'
    
    p_sub = tf.add_paragraph()
    p_sub.text = "Asset Understanding & Response AI"
    p_sub.font.size = Pt(28)
    p_sub.font.bold = True
    p_sub.font.color.rgb = LIGHT_ORANGE
    p_sub.font.name = 'Arial'
    p_sub.space_before = Pt(10)
    
    p_tag = tf.add_paragraph()
    p_tag.text = "The Next-Generation Safety & Operations Copilot for Industrial Plants"
    p_tag.font.size = Pt(18)
    p_tag.font.italic = True
    p_tag.font.color.rgb = GRAY
    p_tag.font.name = 'Arial'
    p_tag.space_before = Pt(20)
    
    # Footer badges
    badge_box = slide1.shapes.add_textbox(Inches(0.75), Inches(6.0), Inches(10), Inches(0.5))
    btf = badge_box.text_frame
    bp = btf.paragraphs[0]
    bp.text = "Enterprise AI  |  Zero-Incident Vision  |  Gemini 1.5 Flash Enabled"
    bp.font.size = Pt(12)
    bp.font.color.rgb = GRAY

    # -------------------------------------------------------------
    # SLIDE 2: The Core Challenge
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide2, NAVY)
    add_header(slide2, "The Core Challenge", "Industrial Plant Operations Are Fragile & Siloed")
    
    # Left Column: Legacy Chaos
    left_box = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(2.0), Inches(5.6), Inches(4.5)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = DARK_NAVY
    left_box.line.color.rgb = RED
    left_box.line.width = Pt(1.5)
    
    ltf = left_box.text_frame
    ltf.word_wrap = True
    ltf.margin_left = ltf.margin_right = Inches(0.3)
    ltf.margin_top = Inches(0.3)
    
    lp = ltf.paragraphs[0]
    lp.text = "THE LEGACY CHAOS"
    lp.font.size = Pt(18)
    lp.font.bold = True
    lp.font.color.rgb = RED
    lp.space_after = Pt(15)
    
    challenges = [
        "Static PDF Manuals: Thousand-page SOPs sit unread in directories while critical choices are made on guess-work.",
        "Disconnected Incident Logs: Lessons from past mechanical failures remain buried in outdated local databases.",
        "Manual Permitting Check: Reviewing Permits to Work (PTWs) takes hours of manual checking, inviting human errors.",
        "Reactive Alarm Alerts: Warnings only trigger after critical machinery shuts down or accidents occur."
    ]
    for ch in challenges:
        p = ltf.add_paragraph()
        p.text = "❌  " + ch
        p.font.size = Pt(12.5)
        p.font.color.rgb = GRAY
        p.space_after = Pt(10)
        
    # Right Column: AURA Paradigm
    right_box = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.98), Inches(2.0), Inches(5.6), Inches(4.5)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = DARK_NAVY
    right_box.line.color.rgb = GREEN
    right_box.line.width = Pt(1.5)
    
    rtf = right_box.text_frame
    rtf.word_wrap = True
    rtf.margin_left = rtf.margin_right = Inches(0.3)
    rtf.margin_top = Inches(0.3)
    
    rp = rtf.paragraphs[0]
    rp.text = "THE AURA PARADIGM"
    rp.font.size = Pt(18)
    rp.font.bold = True
    rp.font.color.rgb = GREEN
    rp.space_after = Pt(15)
    
    solutions = [
        "Instant Semantic Access: Unified context-aware query engine retrieves exact clauses across files in seconds.",
        "Connected Knowledge Graph: Incidents, assets, and safety guidelines are auto-linked dynamically.",
        "Automated Compliance Auditor: AI parses PTWs against SOP criteria to flag compliance gaps instantly.",
        "Proactive Risk Pulse Ticker: Continuous background scanning of active telemetry and log data."
    ]
    for sol in solutions:
        p = rtf.add_paragraph()
        p.text = "✅  " + sol
        p.font.size = Pt(12.5)
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)

    # -------------------------------------------------------------
    # SLIDE 3: The AURA Solution
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide3, NAVY)
    add_header(slide3, "The AURA Solution", "A Centralized Cognitive Hub for Plant Operators")
    
    # 4 Grid blocks
    grid_data = [
        ("💬 Natural Copilot", "Voice and text chat driven by custom RAG, delivering precise answers linked directly to original source citations.", 0.75, 2.0),
        ("🛡️ Compliance Engine", "Immediate validation of dynamic permits. Auto-generates complete Word (.docx) Audit Packs for external compliance.", 6.98, 2.0),
        ("🕸️ D3 Knowledge Graphs", "Interactive force-directed relationships showing how physical assets, safety codes, and maintenance records link together.", 0.75, 4.4),
        ("🚨 Risk Pulse Monitoring", "A scrolling ticker displaying severity-colored alerts about plant equipment failures, hot-work, and environment shifts.", 6.98, 4.4)
    ]
    
    for title, desc, x, y in grid_data:
        box = slide3.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.6), Inches(2.1)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = DARK_NAVY
        box.line.color.rgb = ORANGE
        box.line.width = Pt(1)
        
        btf = box.text_frame
        btf.word_wrap = True
        btf.margin_left = btf.margin_right = Inches(0.25)
        btf.margin_top = Inches(0.2)
        
        bp = btf.paragraphs[0]
        bp.text = title
        bp.font.size = Pt(16)
        bp.font.bold = True
        bp.font.color.rgb = WHITE
        bp.space_after = Pt(8)
        
        bp2 = btf.add_paragraph()
        bp2.text = desc
        bp2.font.size = Pt(12)
        bp2.font.color.rgb = GRAY

    # -------------------------------------------------------------
    # SLIDE 4: Key Capabilities & Features
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide4, NAVY)
    add_header(slide4, "Key Capabilities & Features", "High-Fidelity Contextual Modules")
    
    features = [
        ("RAG COPILOT", "Dynamic Persona-Tailored Retrieval", 
         "AURA adapts its response based on logged-in profiles. A Safety Officer receives safety regulatory codes, a Field Technician gets hands-on steps and hazard warnings, and an Engineer receives technical specifications. Responses include precise citation backlinks to source PDFs.",
         0.75),
        ("SAFETY COMPLIANCE", "Automated Permit Gap Analysis",
         "Ingest permits in real-time. The compliance engine cross-checks isolation coordinates, required protective equipment (PPE), and hot-work conditions against OISD standards, highlighting Critical and Warning compliance mismatches instantly.",
         4.83),
        ("D3 NETWORK ENGINE", "Entity-Relationship Incident Mesh",
         "AURA parses historical logs to build a knowledge network. It automatically links physical machinery (e.g. Boiler #2) to failure modes (e.g. gasket fatigue), safety guides, and operators, allowing visual root-cause analysis (RCA).",
         8.91)
    ]
    
    for title, subtitle, body, x in features:
        card = slide4.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.0), Inches(3.7), Inches(4.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_NAVY
        card.line.color.rgb = RGBColor(40, 50, 70)
        card.line.width = Pt(1)
        
        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_left = ctf.margin_right = Inches(0.25)
        ctf.margin_top = Inches(0.25)
        
        cp1 = ctf.paragraphs[0]
        cp1.text = title
        cp1.font.size = Pt(12)
        cp1.font.bold = True
        cp1.font.color.rgb = LIGHT_ORANGE
        cp1.space_after = Pt(6)
        
        cp2 = ctf.add_paragraph()
        cp2.text = subtitle
        cp2.font.size = Pt(16)
        cp2.font.bold = True
        cp2.font.color.rgb = WHITE
        cp2.space_after = Pt(12)
        
        cp3 = ctf.add_paragraph()
        cp3.text = body
        cp3.font.size = Pt(11.5)
        cp3.font.color.rgb = GRAY

    # -------------------------------------------------------------
    # SLIDE 5: System Architecture
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide5, NAVY)
    add_header(slide5, "System Architecture", "Secure local indexing with intelligent cloud reasoning")
    
    # Box: Data Ingestion
    ingest_box = slide5.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(2.0), Inches(5.6), Inches(3.5)
    )
    ingest_box.fill.solid()
    ingest_box.fill.fore_color.rgb = DARK_NAVY
    ingest_box.line.color.rgb = GREEN
    
    itf = ingest_box.text_frame
    itf.word_wrap = True
    itf.margin_left = itf.margin_right = Inches(0.3)
    itf.margin_top = Inches(0.3)
    
    ip = itf.paragraphs[0]
    ip.text = "DOCUMENT INGESTION WORKFLOW"
    ip.font.size = Pt(16)
    ip.font.bold = True
    ip.font.color.rgb = GREEN
    ip.space_after = Pt(15)
    
    ingest_steps = [
        "1. File Upload: Ingest plant SOPs, PTWs (.pdf, .txt).",
        "2. Parse & Segment: fastapi + pypdf extracts raw paragraphs.",
        "3. Local Vectorization: Embed texts with sentence-transformers.",
        "4. Persistent Indexing: Save indices locally in ChromaDB database."
    ]
    for step in ingest_steps:
        p = itf.add_paragraph()
        p.text = step
        p.font.size = Pt(12.5)
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)
        
    # Box: Query Retrieval
    query_box = slide5.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.98), Inches(2.0), Inches(5.6), Inches(3.5)
    )
    query_box.fill.solid()
    query_box.fill.fore_color.rgb = DARK_NAVY
    query_box.line.color.rgb = ORANGE
    
    qtf = query_box.text_frame
    qtf.word_wrap = True
    qtf.margin_left = qtf.margin_right = Inches(0.3)
    qtf.margin_top = Inches(0.3)
    
    qp = qtf.paragraphs[0]
    qp.text = "USER QUERY / RAG WORKFLOW"
    qp.font.size = Pt(16)
    qp.font.bold = True
    qp.font.color.rgb = ORANGE
    qp.space_after = Pt(15)
    
    query_steps = [
        "1. SSE Stream: React client queries backend endpoint.",
        "2. Persona Search: Retrieve contexts from ChromaDB via similarity.",
        "3. LLM Processing: Send query + context to Gemini 1.5 Flash.",
        "4. Real-time Streams: Tokens & citations streamed back using SSE."
    ]
    for step in query_steps:
        p = qtf.add_paragraph()
        p.text = step
        p.font.size = Pt(12.5)
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)
        
    # Bottom tech metrics
    tech_box = slide5.shapes.add_textbox(Inches(0.75), Inches(5.8), Inches(11.83), Inches(1.0))
    ttf = tech_box.text_frame
    tp = ttf.paragraphs[0]
    tp.text = "Key Architecture Advantages:"
    tp.font.size = Pt(14)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    tp.space_after = Pt(4)
    
    tp2 = ttf.add_paragraph()
    tp2.text = "▪ 1M Tokens Context Window (ingest heavy operations manuals)   ▪ Local Embeddings Core (IP and security isolation)   ▪ SSE Token Streams (immediate response delivery)"
    tp2.font.size = Pt(12)
    tp2.font.color.rgb = GRAY

    # -------------------------------------------------------------
    # SLIDE 6: Market Value & Traction
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide6, NAVY)
    add_header(slide6, "Market Value & Traction", "Empowering Heavy Industries Toward Zero Failure Rates")
    
    metrics = [
        ("90%", "Audit Time Reduction", "Permit reviews trimmed from several hours to instantaneous gap checker audits, automatically generating Word format reports."),
        ("100%", "Compliance Traceability", "Responses are back-linked directly to registered guidelines and manuals with references, eliminating guesswork."),
        ("ZERO", "Safety Violations Target", "Proactive hazard matching blocks scheduling conflicts and safety breaches before field work ever begins.")
    ]
    
    for i, (stat, header, desc) in enumerate(metrics):
        col_x = 0.75 + (i * 4.0)
        m_box = slide6.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(col_x), Inches(2.0), Inches(3.7), Inches(3.2)
        )
        m_box.fill.solid()
        m_box.fill.fore_color.rgb = DARK_NAVY
        m_box.line.color.rgb = ORANGE if i == 2 else RGBColor(40, 50, 70)
        m_box.line.width = Pt(1.5)
        
        mtf = m_box.text_frame
        mtf.word_wrap = True
        mtf.margin_left = mtf.margin_right = Inches(0.25)
        mtf.margin_top = Inches(0.25)
        
        mp1 = mtf.paragraphs[0]
        mp1.text = stat
        mp1.font.size = Pt(36)
        mp1.font.bold = True
        mp1.font.color.rgb = LIGHT_ORANGE
        mp1.space_after = Pt(6)
        
        mp2 = mtf.add_paragraph()
        mp2.text = header
        mp2.font.size = Pt(15)
        mp2.font.bold = True
        mp2.font.color.rgb = WHITE
        mp2.space_after = Pt(10)
        
        mp3 = mtf.add_paragraph()
        mp3.text = desc
        mp3.font.size = Pt(11.5)
        mp3.font.color.rgb = GRAY
        
    # Sector focus
    focus_box = slide6.shapes.add_textbox(Inches(0.75), Inches(5.6), Inches(11.83), Inches(1.2))
    ftf = focus_box.text_frame
    fp = ftf.paragraphs[0]
    fp.text = "TARGET HIGH-HAZARD SECTORS"
    fp.font.size = Pt(12)
    fp.font.bold = True
    fp.font.color.rgb = LIGHT_ORANGE
    fp.space_after = Pt(8)
    
    fp2 = ftf.add_paragraph()
    fp2.text = "🛢️ Oil & Gas Refineries   |   ⚗️ Chemical Processing Plants   |   ⚡ Power Generation Grids   |   🏗️ Heavy Manufacturing"
    fp2.font.size = Pt(14)
    fp2.font.bold = True
    fp2.font.color.rgb = WHITE

    # -------------------------------------------------------------
    # SLIDE 7: Vision & Roadmap
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide7, NAVY)
    add_header(slide7, "Roadmap & Vision", "Unifying Operational Knowledge with Industrial Safety")
    
    # Quote box
    q_box = slide7.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.0), Inches(11.83), Inches(1.3)
    )
    q_box.fill.solid()
    q_box.fill.fore_color.rgb = DARK_NAVY
    q_box.line.color.rgb = ORANGE
    q_box.line.width = Pt(1)
    
    qtf = q_box.text_frame
    qtf.word_wrap = True
    qtf.margin_left = qtf.margin_right = Inches(0.4)
    qtf.margin_top = Inches(0.2)
    
    qp = qtf.paragraphs[0]
    qp.text = '"Our mission is to replace fragmented documentation with dynamic plant intelligence, ensuring every worker goes home safely, every day."'
    qp.font.size = Pt(15)
    qp.font.italic = True
    qp.font.color.rgb = WHITE
    qp.alignment = PP_ALIGN.CENTER
    
    # Timeline
    timeline = [
        ("Phase 1: RAG Copilot", "Seeded vector store, multi-persona chat, citation engine and instant source document lookup.", 0.75),
        ("Phase 2: Auditing & Graph Mesh", "Permit-to-Work safety analysis, D3 interactive entity-relationship maps, automatic Word document audits.", 4.83),
        ("Phase 3: SCADA & Agents", "Connecting real-time sensor streams and PLC inputs, alerting operators of safety hazards automatically.", 8.91)
    ]
    
    for title, desc, x in timeline:
        t_box = slide7.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.7), Inches(3.7), Inches(2.8)
        )
        t_box.fill.solid()
        t_box.fill.fore_color.rgb = DARK_NAVY
        t_box.line.color.rgb = ORANGE if "Phase 3" in title else RGBColor(40, 50, 70)
        t_box.line.width = Pt(1)
        
        ttf = t_box.text_frame
        ttf.word_wrap = True
        ttf.margin_left = ttf.margin_right = Inches(0.25)
        ttf.margin_top = Inches(0.25)
        
        tp1 = ttf.paragraphs[0]
        tp1.text = title.split(":")[0]
        tp1.font.size = Pt(11)
        tp1.font.bold = True
        tp1.font.color.rgb = LIGHT_ORANGE
        tp1.space_after = Pt(4)
        
        tp2 = ttf.add_paragraph()
        tp2.text = title.split(":")[1].strip()
        tp2.font.size = Pt(14)
        tp2.font.bold = True
        tp2.font.color.rgb = WHITE
        tp2.space_after = Pt(10)
        
        tp3 = ttf.add_paragraph()
        tp3.text = desc
        tp3.font.size = Pt(11.5)
        tp3.font.color.rgb = GRAY
        
    # Save the file
    out_path = "C:\\Users\\AASHI JAIN\\OneDrive\\Documents\\AURA\\AURA_Pitch_Deck.pptx"
    prs.save(out_path)
    print(f"[AURA] Pitch Deck presentation created successfully at: {out_path}")

if __name__ == "__main__":
    create_pitch_deck()
